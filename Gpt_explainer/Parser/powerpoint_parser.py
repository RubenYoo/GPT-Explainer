from pptx import Presentation
from collections.abc import Generator
import os


class PowerpointParser:
    """
    This class is used to parse a PowerPoint file and extract the text from each slide
    """
    pptx_file: str

    def __init__(self, file_path: str) -> None:
        """
        Initialize the class
        :param file_path: The path to the PowerPoint file
        """

        if os.path.isfile(file_path):
            self.pptx_file = file_path
        else:
            raise FileNotFoundError("The powerpoint file was not found")

    def extract_text_from_slide(self) -> Generator:
        """
        Extract the text from each slide
        :return: A generator that yields the text from each slide
        """

        # Open the PowerPoint file
        with open(self.pptx_file, "rb") as presentation_file:
            presentation = Presentation(presentation_file)

            # Iterate through each slide and extract the text
            for slide in presentation.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_frames = shape.text_frame
                        slide_text += '\n'.join(
                            [run.text.strip() for paragraph in text_frames.paragraphs for run in paragraph.runs])

                if slide_text:
                    yield slide_text

